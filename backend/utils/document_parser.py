"""
Document parsing utilities for various file formats.
Supports PPTX, DOCX, PDF, and images (with OCR).
"""

from pathlib import Path
from typing import Optional
import tempfile

from pptx import Presentation
from docx import Document as DocxDocument
import PyPDF2
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
from loguru import logger


class DocumentParser:
    """Parser for educational documents."""
    
    @staticmethod
    def parse_pptx(file_path: str) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                text_content.append("\n".join(slide_text))
            
            full_text = "\n\n".join(text_content)
            logger.info(f"Extracted text from {len(prs.slides)} slides in PPTX")
            
            return full_text
            
        except Exception as e:
            logger.error(f"Error parsing PPTX file {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_docx(file_path: str) -> str:
        """Extract text from Word document."""
        try:
            doc = DocxDocument(file_path)
            text_content = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    text_content.append(row_text)
            
            full_text = "\n".join(text_content)
            logger.info(f"Extracted {len(text_content)} paragraphs from DOCX")
            
            return full_text
            
        except Exception as e:
            logger.error(f"Error parsing DOCX file {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_pdf(file_path: str) -> str:
        """Extract text from PDF document."""
        try:
            text_content = []
            
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(f"--- Page {page_num} ---\n{text}")
            
            full_text = "\n\n".join(text_content)
            logger.info(f"Extracted text from {len(pdf_reader.pages)} pages in PDF")
            
            return full_text
            
        except Exception as e:
            logger.error(f"Error parsing PDF file {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_pdf_with_ocr(file_path: str) -> str:
        """
        Extract text from PDF using OCR (for scanned documents).
        Supports Vietnamese language.
        """
        try:
            text_content = []
            
            # Convert PDF to images
            images = convert_from_path(file_path, dpi=300)
            
            for page_num, image in enumerate(images, 1):
                # Perform OCR with Vietnamese language support
                text = pytesseract.image_to_string(image, lang='vie+eng')
                if text.strip():
                    text_content.append(f"--- Page {page_num} ---\n{text}")
            
            full_text = "\n\n".join(text_content)
            logger.info(f"Extracted text from {len(images)} pages using OCR")
            
            return full_text
            
        except Exception as e:
            logger.error(f"Error parsing PDF with OCR {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_image(file_path: str) -> str:
        """
        Extract text from image using OCR.
        Supports Vietnamese language.
        """
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang='vie+eng')
            
            logger.info(f"Extracted text from image: {Path(file_path).name}")
            
            return text.strip()
            
        except Exception as e:
            logger.error(f"Error parsing image {file_path}: {e}")
            raise
    
    @classmethod
    def parse_file(cls, file_path: str, use_ocr: bool = False) -> str:
        """
        Parse file based on extension.
        
        Args:
            file_path: Path to file
            use_ocr: Whether to use OCR for PDFs
            
        Returns:
            Extracted text content
        """
        path = Path(file_path)
        extension = path.suffix.lower()
        
        parsers = {
            ".pptx": cls.parse_pptx,
            ".docx": cls.parse_docx,
            ".pdf": cls.parse_pdf_with_ocr if use_ocr else cls.parse_pdf,
            ".jpg": cls.parse_image,
            ".jpeg": cls.parse_image,
            ".png": cls.parse_image,
        }
        
        parser = parsers.get(extension)
        if parser is None:
            raise ValueError(f"Unsupported file format: {extension}")
        
        return parser(str(file_path))
